"""V15.1 多模态 P3：FileParserRegistry——统一文件解析。

支持 pdf/docx/xlsx/pptx/txt → Markdown。
WHY 统一接口：Agent 不关心文件类型，一个 Tool 全搞定。
"""

from __future__ import annotations

from pathlib import Path

from orbit.tools.models import ToolPermission, ToolSchema

PARSER_SCHEMA = ToolSchema(
    name="file_parser",
    version="1.0.0",
    description="解析文件内容——PDF/Word/Excel/PPT/文本 → Markdown",
    parameters={
        "file_path": {"type": "string", "description": "文件路径"},
    },
    permissions=[ToolPermission.READ],
    allowed_agents=["qa", "developer", "clarifier", "architect", "reviewer", "chatter"],
    timeout_seconds=30,
    is_async=True,
)

# P2-5: PDF 页数上限——防 OOM
MAX_PDF_PAGES = 200
MAX_XLSX_ROWS = 500


async def file_parser(file_path: str) -> dict:
    """统一文件解析入口。

    P2-4 fix: 用 asyncio.to_thread 包装同步 I/O，不阻塞事件循环。
    """
    import asyncio

    # P0 (PR#297): workspace 隔离——防止通过 file_parser 读取工作区外敏感文件
    from orbit.tools.filesystem import _guard_path

    path = _guard_path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"文件不存在: {file_path}")

    ext = path.suffix.lower()

    if ext == ".pdf":
        return await asyncio.to_thread(_parse_pdf, path)
    elif ext == ".docx":
        # P1-1 fix: 移除 .doc（python-docx 不支持二进制 .doc）
        return await asyncio.to_thread(_parse_docx, path)
    elif ext in (".xlsx", ".xls"):
        return await asyncio.to_thread(_parse_xlsx, path)
    elif ext == ".pptx":
        return await asyncio.to_thread(_parse_pptx, path)
    elif ext in (".txt", ".md", ".py", ".js", ".ts", ".vue", ".json", ".yaml", ".yml", ".toml", ".cfg", ".ini"):
        return await asyncio.to_thread(_parse_text, path)
    else:
        raise ValueError(f"不支持的文件类型: {ext}")


def _parse_pdf(path: Path) -> dict:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    total = len(reader.pages)
    pages = []
    limit = min(total, MAX_PDF_PAGES)  # P2-5: 页数上限
    for i, page in enumerate(reader.pages):
        if i >= limit:
            pages.append(f"... (省略 {total - limit} 页)")
            break
        text = page.extract_text()
        if text:
            pages.append(text)
    return {"text": "\n\n".join(pages), "pages": total, "file_type": "pdf"}


def _parse_docx(path: Path) -> dict:
    from docx import Document
    doc = Document(str(path))
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text for cell in row.cells]
            paras.append(" | ".join(cells))
    return {"text": "\n".join(paras), "pages": 1, "file_type": "docx"}


def _parse_xlsx(path: Path) -> dict:
    from openpyxl import load_workbook
    wb = load_workbook(str(path), read_only=True, data_only=True)
    sheets = []
    for name in wb.sheetnames:
        ws = wb[name]
        rows = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= MAX_XLSX_ROWS:
                rows.append(f"... (省略剩余行)")
                break
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append(" | ".join(cells))
        sheets.append(f"## {name}\n" + "\n".join(rows))
    wb.close()
    return {"text": "\n\n".join(sheets), "pages": len(wb.sheetnames), "file_type": "xlsx"}


def _parse_pptx(path: Path) -> dict:
    from pptx import Presentation
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text.strip())
        if texts:
            slides.append(f"## Slide {i+1}\n" + "\n".join(texts))
    return {"text": "\n\n".join(slides), "pages": len(prs.slides), "file_type": "pptx"}


def _parse_text(path: Path) -> dict:
    text = path.read_text(encoding="utf-8", errors="replace")
    return {"text": text, "pages": 1, "file_type": path.suffix.lstrip(".")}


def _register():
    from orbit.tools.registry import get_registry
    get_registry().register(PARSER_SCHEMA, file_parser)

_register()
