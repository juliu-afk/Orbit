"""V15.1 多模态 P3：Office 文件生成 Tool。

/create <type> → 生成 .xlsx / .docx / .pptx 文件。
WHY Agent 需要写文件：测试报告、数据导出、文档生成。
"""

from __future__ import annotations

from pathlib import Path

from orbit.tools.models import ToolPermission, ToolSchema

# P1-4: 路径白名单——防止任意文件写入
_ALLOWED_DIRS = ("./output", "/tmp", "./data", "./exports")  # P2-A: 移除 "." 防源码覆盖


def _validate_output_path(output_path: str) -> Path:
    """验证输出路径在允许的目录内。"""
    path = Path(output_path).resolve()
    allowed = [Path(d).resolve() for d in _ALLOWED_DIRS]
    for d in allowed:
        try:
            path.relative_to(d)
            return path
        except ValueError:
            continue
    raise ValueError(f"禁止写入路径: {output_path}。允许的目录: {_ALLOWED_DIRS}")


OFFICE_SCHEMA = ToolSchema(
    name="office_create",
    version="1.0.0",
    description="生成 Office 文件——Excel/Word/PPT",
    parameters={
        "file_type": {"type": "string", "description": "xlsx | docx | pptx"},
        "output_path": {"type": "string", "description": "输出文件路径"},
        "content": {"type": "object", "description": "内容——xlsx: [[row1],[row2],...], docx: [paragraph,...], pptx: [{title, bullets},...]"},
    },
    permissions=[ToolPermission.WRITE],
    allowed_agents=["developer", "qa", "architect"],
    timeout_seconds=30,
    is_async=True,
)


async def office_create(file_type: str, output_path: str, content: dict | list) -> dict:
    """生成 Office 文件。

    P1-4 fix: 路径白名单验证。
    P2-4 fix: asyncio.to_thread 包装同步 I/O。
    """
    import asyncio

    path = _validate_output_path(output_path)

    if file_type == "xlsx":
        return await asyncio.to_thread(_create_xlsx, path, content)
    elif file_type == "docx":
        return await asyncio.to_thread(_create_docx, path, content)
    elif file_type == "pptx":
        return await asyncio.to_thread(_create_pptx, path, content)
    else:
        raise ValueError(f"不支持的文件类型: {file_type}，支持: xlsx/docx/pptx")


def _create_xlsx(path: Path, content: dict | list) -> dict:
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"

    rows = content if isinstance(content, list) else content.get("rows", [])
    for row in rows:
        ws.append(row if isinstance(row, list) else [row])

    wb.save(str(path))
    return {"path": str(path), "file_type": "xlsx", "size_bytes": path.stat().st_size}


def _create_docx(path: Path, content: dict | list) -> dict:
    from docx import Document
    doc = Document()

    title = content.get("title", "") if isinstance(content, dict) else ""
    if title:
        doc.add_heading(title, level=1)

    paragraphs = content if isinstance(content, list) else content.get("paragraphs", [])
    for p in paragraphs:
        doc.add_paragraph(str(p))

    doc.save(str(path))
    return {"path": str(path), "file_type": "docx", "size_bytes": path.stat().st_size}


def _create_pptx(path: Path, content: dict | list) -> dict:
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()

    slides = content if isinstance(content, list) else content.get("slides", [])
    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if isinstance(slide_data, dict):
            slide.shapes.title.text = slide_data.get("title", "")
            bullets = slide_data.get("bullets", [])
            if bullets and slide.placeholders[1].has_text_frame:
                tf = slide.placeholders[1].text_frame
                for b in bullets:
                    p = tf.add_paragraph()
                    p.text = str(b)
                    p.level = 0

    prs.save(str(path))
    return {"path": str(path), "file_type": "pptx", "size_bytes": path.stat().st_size}


def _register():
    from orbit.tools.registry import get_registry
    get_registry().register(OFFICE_SCHEMA, office_create)

_register()
